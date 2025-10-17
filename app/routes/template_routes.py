import os
import re
from fastapi import APIRouter, UploadFile, File, HTTPException
from fastapi.responses import JSONResponse
from app.session_data import session_data


router = APIRouter();

UPLOAD_DIR = "app/static/templates"
os.makedirs(UPLOAD_DIR, exist_ok=True)

PLACEHOLDER_PATTERN = r"<<(.*?)>>"

@router.post("/upload-template")
async def upload_template(file: UploadFile = File(...)):
    allowed_extensions = ["pdf", "ppt", "pptx"]
    ext = file.filename.split(".")[-1].lower()

    if ext not in allowed_extensions:
        raise HTTPException(
            status_code=400,
            detail="Invalid file type. Allowed: pdf, ppt, pptx"
        )

    file_path = os.path.join(UPLOAD_DIR, file.filename)

    with open(file_path, "wb") as f:
        f.write(await file.read())

    # ✅ Save to session_data
    session_data["template_file"] = file.filename

    placeholders = []
    if ext == "pdf":
        placeholders = extract_text_placeholders(file_path)
    elif ext in ["ppt", "pptx"]:
        placeholders = extract_ppt_placeholders(file_path)

    return JSONResponse({
        "message": "Template uploaded successfully",
        "fileName": file.filename,
        "placeholders": placeholders
    })



def extract_text_placeholders(file_path: str):
    import fitz  # PyMuPDF
    placeholders = set()
    doc = fitz.open(file_path)

    for page in doc:
        text = page.get_text()
        matches = re.findall(PLACEHOLDER_PATTERN, text)
        for match in matches:
            placeholders.add(match.strip())

    doc.close()
    return list(placeholders)


def extract_ppt_placeholders(file_path: str):
    from pptx import Presentation
    placeholders = set()
    prs = Presentation(file_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                matches = re.findall(PLACEHOLDER_PATTERN, shape.text)
                for match in matches:
                    placeholders.add(match.strip())
    return list(placeholders)